"""PDF Studio backend - FastAPI + PyMuPDF."""
import hashlib
import json
import logging
import os
import re
import secrets
import sys
import urllib.error
import urllib.request
from pathlib import Path
from typing import Any

# Load environment variables from .env file if it exists
ENV_PATH = Path(__file__).parent / ".env"
if ENV_PATH.exists():
    with open(ENV_PATH, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                key, val = line.split("=", 1)
                os.environ[key.strip()] = val.strip()

from fastapi import FastAPI, File, Form, HTTPException, UploadFile, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, Response
from fastapi.exceptions import RequestValidationError
from starlette.exceptions import HTTPException as StarletteHTTPException
from pydantic import BaseModel

import auth
import pdf_ops
import storage
import fonts
from compare_engine.engine import ComparisonEngine
from compare_engine.report import ComparisonReportGenerator

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)]
)
logger = logging.getLogger(__name__)

app = FastAPI(title="PDF Studio API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000", "app://localhost"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["*"],
)


@app.exception_handler(StarletteHTTPException)
async def http_exception_handler(request: Request, exc: StarletteHTTPException):
    return JSONResponse(
        status_code=exc.status_code,
        content={"error": exc.detail, "detail": exc.detail},
    )


@app.exception_handler(RequestValidationError)
async def validation_exception_handler(request: Request, exc: RequestValidationError):
    errors = []
    for err in exc.errors():
        loc = " -> ".join(str(x) for x in err.get("loc", []))
        msg = err.get("msg")
        errors.append(f"{loc}: {msg}")
    error_msg = "; ".join(errors) or "Validation error"
    return JSONResponse(
        status_code=422,
        content={"error": error_msg, "detail": error_msg},
    )


def _meta_or_404(file_id: str) -> dict:
    meta = storage.get(file_id)
    if not meta:
        raise HTTPException(404, "File not found")
    if not storage.path_for(file_id).exists():
        raise HTTPException(410, "File data is missing")
    return meta


def _safe_name(name: str) -> str:
    name = re.sub(r"[^\w\-. ()\[\]]+", "_", name).strip() or "document.pdf"
    if not name.lower().endswith(".pdf"):
        name += ".pdf"
    return name


def _stem(meta: dict) -> str:
    return re.sub(r"\.pdf$", "", meta["name"], flags=re.I)


# ------------------------------------------------------------ authentication

class RegisterBody(BaseModel):
    name: str
    email: str
    password: str


class LoginBody(BaseModel):
    email: str
    password: str


class GoogleLoginBody(BaseModel):
    access_token: str


@app.post("/api/auth/register")
def register(body: RegisterBody, response: Response):
    name = body.name.strip()
    email = body.email.strip().lower()
    password = body.password

    if not name or not email or not password:
        raise HTTPException(400, "Name, email, and password are required")

    if len(password) < 6:
        raise HTTPException(400, "Password must be at least 6 characters long")

    existing_user = auth.get_user_by_email(email)
    if existing_user:
        raise HTTPException(400, "User already exists with this email")

    user = auth.create_user(name, email, password)
    token = auth.sign_jwt({"userId": user["id"], "email": user["email"], "name": user["name"]})

    response.set_cookie(
        key="qt_token",
        value=token,
        httponly=True,
        max_age=24 * 60 * 60,  # 24 hours
        samesite="lax",
        path="/"
    )

    return {
        "success": True,
        "user": {
            "id": user["id"],
            "email": user["email"],
            "name": user["name"]
        }
    }


@app.post("/api/auth/login")
def login(body: LoginBody, response: Response):
    email = body.email.strip().lower()
    password = body.password

    if not email or not password:
        raise HTTPException(400, "Email and password are required")

    user = auth.get_user_by_email(email)
    if not user:
        raise HTTPException(401, "Invalid email or password")

    if not auth.verify_password(password, user["passwordHash"]):
        raise HTTPException(401, "Invalid email or password")

    token = auth.sign_jwt({"userId": user["id"], "email": user["email"], "name": user["name"]})

    response.set_cookie(
        key="qt_token",
        value=token,
        httponly=True,
        max_age=24 * 60 * 60,  # 24 hours
        samesite="lax",
        path="/"
    )

    return {
        "success": True,
        "user": {
            "id": user["id"],
            "email": user["email"],
            "name": user["name"]
        }
    }


@app.post("/api/auth/google")
def google_login(body: GoogleLoginBody, response: Response):
    access_token = body.access_token.strip()
    if not access_token:
        raise HTTPException(400, "Google access token is required")

    req = urllib.request.Request(
        "https://www.googleapis.com/oauth2/v3/userinfo",
        headers={"Authorization": f"Bearer {access_token}"}
    )
    try:
        with urllib.request.urlopen(req) as r:
            resp_data = json.loads(r.read().decode("utf-8"))
    except urllib.error.HTTPError as e:
        err_msg = e.read().decode("utf-8")
        raise HTTPException(401, f"Failed to verify Google token: {err_msg}")
    except Exception as e:
        raise HTTPException(401, f"Failed to verify Google token: {e}")

    email = resp_data.get("email", "").strip().lower()
    name = resp_data.get("name", "").strip()
    if not email:
        raise HTTPException(401, "Google token did not contain an email")

    user = auth.get_user_by_email(email)
    if not user:
        random_password = secrets.token_hex(16)
        user = auth.create_user(name or email, email, random_password)

    token = auth.sign_jwt({"userId": user["id"], "email": user["email"], "name": user["name"]})

    response.set_cookie(
        key="qt_token",
        value=token,
        httponly=True,
        max_age=24 * 60 * 60,
        samesite="lax",
        path="/"
    )

    return {
        "success": True,
        "user": {
            "id": user["id"],
            "email": user["email"],
            "name": user["name"]
        }
    }


@app.post("/api/auth/logout")
def logout(response: Response):
    response.delete_cookie(key="qt_token", path="/")
    return {"success": True}


@app.get("/api/auth/me")
def me(request: Request):
    token = request.cookies.get("qt_token")
    if not token:
        return {"loggedIn": False}
    payload = auth.verify_jwt(token)
    if not payload:
        return {"loggedIn": False}
    return {
        "loggedIn": True,
        "user": {
            "email": payload.get("email"),
            "name": payload.get("name")
        }
    }


# ------------------------------------------------------------ fonts

@app.get("/api/fonts")
def get_fonts():
    return {"families": fonts.families()}


# ------------------------------------------------------------ files

@app.get("/api/health")
def health():
    return {"ok": True}


@app.get("/health")
def root_health():
    return {"status": "ok", "message": "Backend is running"}


@app.post("/api/files/upload")
async def upload(file: UploadFile = File(...)):
    data = await file.read()
    try:
        pages = pdf_ops.page_count(data)
    except Exception as e:
        logger.error(f"Failed to open PDF {file.filename}: {e}")
        raise HTTPException(400, "Could not open this PDF")
    meta = storage.create(_safe_name(file.filename or "document.pdf"), data, pages)
    return storage.public(meta)


@app.get("/api/files")
def list_files(include_deleted: bool = False):
    return [storage.public(m) for m in storage.list_files(include_deleted=include_deleted)]


@app.get("/api/files/{file_id}")
def get_meta(file_id: str):
    return storage.public(_meta_or_404(file_id))


@app.post("/api/files/{file_id}/open")
def mark_opened(file_id: str):
    _meta_or_404(file_id)
    storage.touch_opened(file_id)
    return {"ok": True}


class RenameBody(BaseModel):
    name: str


@app.patch("/api/files/{file_id}")
def rename(file_id: str, body: RenameBody):
    _meta_or_404(file_id)
    return storage.public(storage.update(file_id, name=_safe_name(body.name)))


@app.api_route("/api/files/{file_id}/content", methods=["GET", "HEAD"])
def content(file_id: str, v: int | None = None):
    meta = _meta_or_404(file_id)
    data = None
    if v is not None and v != meta.get("version", 1):
        found_snap = None
        for snap_name in meta.get("undo", []) + meta.get("redo", []):
            if snap_name.startswith(f"{file_id}-{v}-"):
                found_snap = snap_name
                break
        if found_snap:
            snap_path = storage.HISTORY_DIR / found_snap
            if snap_path.exists():
                data = snap_path.read_bytes()
    if data is None:
        data = storage.path_for(file_id).read_bytes()
    return Response(data, media_type="application/pdf",
                    headers={"Cache-Control": "no-store"})


@app.get("/api/files/{file_id}/download")
def download(file_id: str):
    meta = _meta_or_404(file_id)
    data = storage.path_for(file_id).read_bytes()
    return Response(data, media_type="application/pdf", headers={
        "Content-Disposition": f'attachment; filename="{meta["name"]}"'})


@app.delete("/api/files/{file_id}")
def delete(file_id: str, permanent: bool = False):
    if not storage.get(file_id):
        raise HTTPException(404, "File not found")
    if permanent:
        storage.hard_delete(file_id)
    else:
        storage.soft_delete(file_id)
    return {"ok": True}


@app.post("/api/files/{file_id}/restore")
def restore(file_id: str):
    if not storage.restore(file_id):
        raise HTTPException(404, "File not found")
    return {"ok": True}


# ------------------------------------------------------------ document info

@app.get("/api/files/{file_id}/outline")
def outline(file_id: str):
    _meta_or_404(file_id)
    return pdf_ops.get_outline(storage.path_for(file_id))


@app.get("/api/files/{file_id}/paragraphs")
def paragraphs(file_id: str, page: int = 0):
    _meta_or_404(file_id)
    return pdf_ops.get_paragraphs(storage.path_for(file_id), page)


@app.get("/api/files/{file_id}/notes")
def notes(file_id: str):
    _meta_or_404(file_id)
    return pdf_ops.list_notes(storage.path_for(file_id))


class NoteBody(BaseModel):
    text: str


@app.patch("/api/files/{file_id}/notes/{xref}")
def edit_note(file_id: str, xref: int, body: NoteBody):
    _meta_or_404(file_id)
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        if not pdf_ops.update_note(storage.path_for(file_id), xref, body.text):
            raise HTTPException(404, "Note not found")
        return storage.public(storage.update(file_id, bump_version=True))


@app.delete("/api/files/{file_id}/notes/{xref}")
def remove_note(file_id: str, xref: int):
    _meta_or_404(file_id)
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        if not pdf_ops.delete_note(storage.path_for(file_id), xref):
            raise HTTPException(404, "Note not found")
        return storage.public(storage.update(file_id, bump_version=True))


# ------------------------------------------------------------ form fields

@app.get("/api/files/{file_id}/fields")
def form_fields(file_id: str):
    _meta_or_404(file_id)
    return pdf_ops.list_form_fields(storage.path_for(file_id))


class FieldValue(BaseModel):
    xref: int
    value: str | bool | None = None


class FieldsBody(BaseModel):
    values: list[FieldValue]


@app.post("/api/files/{file_id}/fields")
def set_form_fields(file_id: str, body: FieldsBody):
    _meta_or_404(file_id)
    if not body.values:
        raise HTTPException(400, "No field values supplied")
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        n = pdf_ops.set_form_fields(
            storage.path_for(file_id), {f.xref: f.value for f in body.values})
        if not n:
            raise HTTPException(404, "No matching form fields in this document")
        return storage.public(storage.update(file_id, bump_version=True))


# ------------------------------------------------------------ save edits

class AnnotationsBody(BaseModel):
    annotations: list[dict[str, Any]]


@app.post("/api/files/{file_id}/annotations")
def save_annotations(file_id: str, body: AnnotationsBody):
    _meta_or_404(file_id)
    if not body.annotations:
        raise HTTPException(400, "No edits supplied")
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        try:
            warnings, changed_pages = pdf_ops.bake_annotations(
                storage.path_for(file_id), body.annotations)
        except ValueError as exc:
            raise HTTPException(409, str(exc))
        except Exception as exc:
            raise HTTPException(422, f"Could not apply edits: {exc}")
        meta = storage.update(file_id, bump_version=True,
                              pages=pdf_ops.page_count(
                                  storage.path_for(file_id).read_bytes()))
    out = storage.public(meta)
    out["warnings"] = warnings
    out["changed_pages"] = changed_pages
    return out


@app.post("/api/files/{file_id}/undo")
def undo(file_id: str):
    _meta_or_404(file_id)
    with storage.op_lock(file_id):
        meta = storage.undo(file_id)
        if meta is None:
            raise HTTPException(409, "Nothing to undo")
        pages = pdf_ops.page_count(storage.path_for(file_id).read_bytes())
        return storage.public(storage.update(file_id, pages=pages))


@app.post("/api/files/{file_id}/redo")
def redo(file_id: str):
    _meta_or_404(file_id)
    with storage.op_lock(file_id):
        meta = storage.redo(file_id)
        if meta is None:
            raise HTTPException(409, "Nothing to redo")
        pages = pdf_ops.page_count(storage.path_for(file_id).read_bytes())
        return storage.public(storage.update(file_id, pages=pages))


class WatermarkBody(BaseModel):
    text: str = "CONFIDENTIAL"
    color: str = "#c0392b"
    opacity: float = 0.18
    fontSize: float | None = None
    rotate: float = -45
    pages: str = "all"


@app.post("/api/files/{file_id}/watermark")
def watermark(file_id: str, body: WatermarkBody):
    meta = _meta_or_404(file_id)
    try:
        pages = pdf_ops.parse_ranges(body.pages, meta["pages"])
    except Exception as exc:
        raise HTTPException(400, f"Bad page range: {exc}")
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        pdf_ops.add_watermark(
            storage.path_for(file_id), text=body.text, color=body.color,
            opacity=max(0.02, min(1.0, body.opacity)),
            font_size=body.fontSize, rotate=body.rotate, pages=pages)
        return storage.public(storage.update(file_id, bump_version=True))


# ------------------------------------------------------------ page operations

class PagesBody(BaseModel):
    pages: list[int]


class RotateBody(PagesBody):
    degrees: int = 90


class OrderBody(BaseModel):
    order: list[int]


class ExtractBody(PagesBody):
    name: str | None = None


def _bumped(file_id: str) -> dict:
    path = storage.path_for(file_id)
    return storage.public(storage.update(
        file_id, bump_version=True,
        pages=pdf_ops.page_count(path.read_bytes())))


@app.post("/api/files/{file_id}/pages/rotate")
def rotate(file_id: str, body: RotateBody):
    _meta_or_404(file_id)
    if body.degrees % 90 != 0:
        raise HTTPException(400, "Rotation must be a multiple of 90")
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        pdf_ops.rotate_pages(storage.path_for(file_id), body.pages, body.degrees)
        return _bumped(file_id)


@app.post("/api/files/{file_id}/pages/delete")
def remove_pages(file_id: str, body: PagesBody):
    _meta_or_404(file_id)
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        try:
            pdf_ops.delete_pages(storage.path_for(file_id), body.pages)
        except ValueError as exc:
            raise HTTPException(400, str(exc))
        return _bumped(file_id)


class InsertPageBody(BaseModel):
    after: int = -1  # 0-based page the blank page goes after; -1 = at the front


@app.post("/api/files/{file_id}/pages/insert")
def insert_page(file_id: str, body: InsertPageBody):
    meta = _meta_or_404(file_id)
    if not (-1 <= body.after < meta["pages"]):
        raise HTTPException(400, "Insert position is out of range")
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        pdf_ops.insert_blank_page(storage.path_for(file_id), body.after)
        return _bumped(file_id)


@app.post("/api/files/{file_id}/pages/duplicate")
def duplicate(file_id: str, body: PagesBody):
    _meta_or_404(file_id)
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        pdf_ops.duplicate_pages(storage.path_for(file_id), body.pages)
        return _bumped(file_id)


@app.post("/api/files/{file_id}/pages/reorder")
def reorder(file_id: str, body: OrderBody):
    _meta_or_404(file_id)
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        try:
            pdf_ops.reorder_pages(storage.path_for(file_id), body.order)
        except ValueError as exc:
            raise HTTPException(400, str(exc))
        return _bumped(file_id)


@app.post("/api/files/{file_id}/pages/extract")
def extract(file_id: str, body: ExtractBody):
    meta = _meta_or_404(file_id)
    try:
        data = pdf_ops.extract_pages(storage.path_for(file_id), body.pages)
    except ValueError as exc:
        raise HTTPException(400, str(exc))
    name = _safe_name(body.name or f"{_stem(meta)} (extract).pdf")
    return storage.public(storage.create(name, data, pdf_ops.page_count(data)))


class SplitBody(BaseModel):
    ranges: str  # e.g. "1-3, 4-6, 7-"


@app.post("/api/files/{file_id}/split")
def split(file_id: str, body: SplitBody):
    meta = _meta_or_404(file_id)
    try:
        groups = pdf_ops.parse_range_groups(body.ranges, meta["pages"])
        parts = pdf_ops.split_document(storage.path_for(file_id), groups)
    except ValueError as exc:
        raise HTTPException(400, str(exc))
    created = []
    for i, data in enumerate(parts):
        name = _safe_name(f"{_stem(meta)} (part {i + 1}).pdf")
        created.append(storage.public(
            storage.create(name, data, pdf_ops.page_count(data))))
    return created


class MergeBody(BaseModel):
    file_ids: list[str]
    name: str | None = None


@app.post("/api/files/merge")
def merge(body: MergeBody):
    if len(body.file_ids) < 2:
        raise HTTPException(400, "Pick at least two documents to merge")
    paths = []
    for fid in body.file_ids:
        _meta_or_404(fid)
        paths.append(storage.path_for(fid))
    data = pdf_ops.merge_documents(paths)
    name = _safe_name(body.name or "Merged document.pdf")
    return storage.public(storage.create(name, data, pdf_ops.page_count(data)))


# ------------------------------------------------------------ compare PDF

@app.post("/api/files/compare")
async def compare_files(
    original: UploadFile | None = File(None),
    revised: UploadFile | None = File(None),
    original_id: str | None = Form(None),
    compare_versions: str | None = Form(None)
):
    if compare_versions == "true":
        if not original_id:
            raise HTTPException(400, "original_id must be provided when comparing versions")
        
        meta = storage.get(original_id)
        if not meta:
            raise HTTPException(404, "Document not found")
            
        undo_history = meta.get("undo", [])
        if undo_history:
            path_orig = storage.HISTORY_DIR / undo_history[0]
            if not path_orig.exists():
                logger.warning(f"Original version snapshot {undo_history[0]} not found on disk, falling back to current file")
                path_orig = storage.path_for(original_id)
        else:
            path_orig = storage.path_for(original_id)
            
        if not path_orig.exists():
            raise HTTPException(410, "Original version data is missing")

        path_rev = storage.path_for(original_id)
        if not path_rev.exists():
            raise HTTPException(410, "Current document data is missing")
            
        orig_meta = meta.copy()
        orig_meta["version"] = 1
        if undo_history and (storage.HISTORY_DIR / undo_history[0]).exists():
            orig_meta["size"] = (storage.HISTORY_DIR / undo_history[0]).stat().st_size
            try:
                import fitz
                with fitz.open(str(storage.HISTORY_DIR / undo_history[0])) as doc:
                    orig_meta["pages"] = doc.page_count
            except Exception:
                pass
        
        rev_meta = meta
    else:
        if original_id:
            orig_meta = storage.get(original_id)
            if not orig_meta:
                raise HTTPException(404, "Original document not found on server")
            path_orig = storage.path_for(original_id)
            if not path_orig.exists():
                raise HTTPException(410, "Original document data is missing")
        else:
            if not original:
                raise HTTPException(400, "Either original file or original_id must be provided")
            orig_data = await original.read()
            try:
                pdf_ops.page_count(orig_data)
            except Exception as e:
                logger.error(f"Original file {original.filename} is corrupted: {e}")
                raise HTTPException(400, "Original file is corrupted or not a valid PDF")
            orig_meta = storage.create(
                _safe_name(original.filename or "original.pdf"),
                orig_data,
                pdf_ops.page_count(orig_data)
            )
            path_orig = storage.path_for(orig_meta["id"])

        if not revised:
            raise HTTPException(400, "Revised file must be provided")
        rev_data = await revised.read()
        try:
            pdf_ops.page_count(rev_data)
        except Exception as e:
            logger.error(f"Revised file {revised.filename} is corrupted: {e}")
            raise HTTPException(400, "Revised file is corrupted or not a valid PDF")

        rev_meta = storage.create(
            _safe_name(revised.filename or "revised.pdf"),
            rev_data,
            pdf_ops.page_count(rev_data)
        )
        path_rev = storage.path_for(rev_meta["id"])

    try:
        diff_result = ComparisonEngine.compare_documents(str(path_orig), str(path_rev))
        diff_result["file_id_original"] = orig_meta["id"]
        diff_result["file_id_revised"] = rev_meta["id"]
        diff_result["meta_original"] = storage.public(orig_meta)
        diff_result["meta_revised"] = storage.public(rev_meta)
        return diff_result
    except ValueError as exc:
        raise HTTPException(400, str(exc))
    except Exception as exc:
        logger.error(f"Critical comparison failure: {exc}", exc_info=True)
        raise HTTPException(500, f"Comparison failed: {exc}")


class CompareReportBody(BaseModel):
    summary: dict


@app.post("/api/files/compare/report")
def export_compare_report(body: CompareReportBody):
    try:
        report_data = ComparisonReportGenerator.generate(body.summary)
        return Response(
            report_data,
            media_type="application/pdf",
            headers={"Content-Disposition": 'attachment; filename="Comparison_Report.pdf"'}
        )
    except Exception as exc:
        logger.error(f"Failed to generate report: {exc}", exc_info=True)
        raise HTTPException(500, f"Failed to generate report: {exc}")


# ------------------------------------------------------------ protect PDF

class ProtectBody(BaseModel):
    password: str
    owner_password: str | None = None
    permissions: list[str] | None = None


@app.post("/api/files/{file_id}/protect")
def protect_pdf(file_id: str, body: ProtectBody):
    meta = _meta_or_404(file_id)
    if not body.password:
        raise HTTPException(400, "Password is required")
    with storage.op_lock(file_id):
        storage.snapshot_before_change(file_id)
        try:
            pdf_ops.protect(
                storage.path_for(file_id),
                user_password=body.password,
                owner_password=body.owner_password or body.password,
            )
        except Exception as exc:
            raise HTTPException(500, f"Failed to protect PDF: {exc}")
        return storage.public(storage.update(file_id, bump_version=True))


# ------------------------------------------------------------ compress PDF

@app.post("/api/files/{file_id}/compress")
def compress_pdf_file(file_id: str, body: dict = None):
    """Compress the PDF file size by resampling images and optimizing streams."""
    meta = _meta_or_404(file_id)
    quality = 70
    target_size = None
    if body:
        if "quality" in body:
            quality = int(body["quality"])
        if "target_size" in body and body["target_size"] is not None:
            target_size = int(body["target_size"])

    with storage.op_lock(file_id):
        pdf_path = storage.path_for(file_id)
        data = pdf_path.read_bytes()

        try:
            if target_size is not None and target_size > 0:
                compressed_data = pdf_ops.compress_to_target_size(data, target_size=target_size)
            else:
                compressed_data = pdf_ops.compress(data, quality=quality)
        except Exception as exc:
            logger.error(f"Compression failed: {exc}", exc_info=True)
            raise HTTPException(500, "Unable to compress this PDF. The file may be corrupt or already fully optimized.")

        storage.snapshot_before_change(file_id)
        pdf_path.write_bytes(compressed_data)
        updated_meta = storage.update(file_id, bump_version=True)

    return storage.public(updated_meta)


# ------------------------------------------------------------ summarize PDF

@app.post("/api/files/{file_id}/summarize")
def summarize_pdf_file(file_id: str, body: dict = None):
    """Summarize the PDF file content and return clean Markdown."""
    meta = _meta_or_404(file_id)
    mode = "medium"
    if body and "mode" in body:
        mode = body["mode"]
    if mode not in ("short", "medium", "detailed"):
        mode = "medium"
        
    pdf_path = storage.path_for(file_id)
    if not pdf_path.exists():
        raise HTTPException(404, "PDF file not found on disk.")

    try:
        summary_md = pdf_ops.summarize_pdf(pdf_path, mode=mode)
        if not summary_md or not summary_md.strip():
            return {"summary": "# Summary\n\nCould not generate a summary. The document may not contain extractable text."}
        return {"summary": summary_md}
    except Exception as exc:
        logger.error(f"Summarization failed: {exc}", exc_info=True)
        # Return a graceful response instead of 500
        return {"summary": f"# Summary Generation Error\n\nWe encountered an issue while analyzing your document:\n\n> {str(exc)}\n\nPlease try again or use a different summary mode."}


@app.post("/api/files/summarize/export")
def export_summary(body: dict = None, format: str = "pdf"):
    """Convert summary markdown to PDF or Word (DOCX) and return it as a binary file."""
    import tempfile
    
    if not body or "markdown" not in body or "filename" not in body:
        raise HTTPException(400, "Missing 'markdown' or 'filename' in request body")
        
    markdown = body["markdown"]
    filename = body["filename"]
    
    fmt = format.lower().strip()
    if fmt not in ("pdf", "docx"):
        raise HTTPException(400, "Supported formats: pdf, docx")
        
    safe_filename = re.sub(r'[^\w\s\-\.]', '', filename).strip()
    if not safe_filename:
        safe_filename = "summary"
        
    # Ensure correct extension
    if fmt == "pdf":
        if not safe_filename.endswith(".pdf"):
            safe_filename += ".pdf"
    elif fmt == "docx":
        if not safe_filename.endswith(".docx"):
            safe_filename += ".docx"
            
    # Create temp file
    temp_dir = tempfile.gettempdir()
    temp_path = os.path.join(temp_dir, f"{secrets.token_hex(8)}_{safe_filename}")
    
    try:
        if fmt == "pdf":
            pdf_ops.markdown_to_pdf(markdown, temp_path, safe_filename)
            media_type = "application/pdf"
        else:
            pdf_ops.markdown_to_docx(markdown, temp_path, safe_filename)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            
        with open(temp_path, "rb") as f:
            file_data = f.read()
            
        return Response(
            content=file_data,
            media_type=media_type,
            headers={
                "Content-Disposition": f"attachment; filename=\"{safe_filename}\"",
                "Access-Control-Expose-Headers": "Content-Disposition"
            }
        )
    except Exception as exc:
        logger.error(f"Failed to export summary: {exc}", exc_info=True)
        raise HTTPException(500, f"Failed to export summary: {str(exc)}")
    finally:
        if os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception:
                pass



# ------------------------------------------------------------ export (PPTX hot-reload trigger)

@app.get("/api/files/{file_id}/export")
def export(file_id: str, format: str = "pdf", pages: str = "all", dpi: int = 150):
    try:
        meta = _meta_or_404(file_id)
        fmt = format.lower()
        dpi = max(72, min(600, dpi))
        try:
            page_list = pdf_ops.parse_ranges(pages, meta["pages"])
        except Exception as exc:
            raise HTTPException(400, f"Bad page range: {exc}")
        stem = _stem(meta)

        if fmt == "pdf":
            data = pdf_ops.extract_pages(storage.path_for(file_id), page_list)
            return Response(data, media_type="application/pdf", headers={
                "Content-Disposition": f'attachment; filename="{stem} (export).pdf"'})

        if fmt == "docx":
            # Extract page subset first
            data_pdf = pdf_ops.extract_pages(storage.path_for(file_id), page_list)
            import tempfile
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
                tmp_pdf.write(data_pdf)
                tmp_pdf_path = tmp_pdf.name
            tmp_docx_path = tmp_pdf_path.replace(".pdf", ".docx")
            try:
                from pdf2docx import Converter
                cv = Converter(tmp_pdf_path)
                cv.convert(tmp_docx_path)
                cv.close()
                with open(tmp_docx_path, "rb") as f:
                    docx_bytes = f.read()
            except Exception as exc:
                raise HTTPException(500, f"PDF to Word conversion failed: {exc}")
            finally:
                import os as _os
                for p in [tmp_pdf_path, tmp_docx_path]:
                    try:
                        _os.unlink(p)
                    except Exception:
                        pass
            return Response(docx_bytes, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", headers={
                "Content-Disposition": f'attachment; filename="{stem} (export).docx"'})

        if fmt == "pptx":
            from pptx import Presentation
            from pptx.util import Inches
            import io

            # Render PDF pages to high quality images
            entries = pdf_ops.export_images(storage.path_for(file_id), page_list, fmt="png", dpi=150)
            if not entries:
                raise HTTPException(400, "No valid pages selected")

            # Open PDF to get first page dimensions
            import fitz
            doc = fitz.open(str(storage.path_for(file_id)))
            try:
                first_page = doc[page_list[0]]
                page_w = first_page.rect.width
                page_h = first_page.rect.height
            except Exception:
                page_w, page_h = 612, 792 # Default letter size
            finally:
                doc.close()

            prs = Presentation()
            prs.slide_width = Inches(page_w / 72.0)
            prs.slide_height = Inches(page_h / 72.0)
            blank_slide_layout = prs.slide_layouts[6]

            for name, img_data in entries:
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(io.BytesIO(img_data), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

            ppt_buf = io.BytesIO()
            prs.save(ppt_buf)
            pptx_bytes = ppt_buf.getvalue()

            return Response(pptx_bytes, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={
                "Content-Disposition": f'attachment; filename="{stem} (export).pptx"'})

        if fmt not in ("png", "jpg", "jpeg"):
            raise HTTPException(400, "format must be pdf, docx, pptx, png or jpg")
        fmt = "png" if fmt == "png" else "jpg"
        entries = pdf_ops.export_images(storage.path_for(file_id), page_list, fmt, dpi)
        if not entries:
            raise HTTPException(400, "No valid pages selected")
        if len(entries) == 1:
            name, data = entries[0]
            media = "image/png" if fmt == "png" else "image/jpeg"
            return Response(data, media_type=media, headers={
                "Content-Disposition": f'attachment; filename="{stem} {name}"'})
        data = pdf_ops.zip_bytes(entries)
        return Response(data, media_type="application/zip", headers={
            "Content-Disposition": f'attachment; filename="{stem} ({fmt} export).zip"'})
    except Exception as exc:
        logger.error(f"Export failed: {exc}", exc_info=True)
        raise HTTPException(500, f"Export failed: {exc}")


# ------------------------------------------------------------ PDF to Word

@app.post("/api/convert/pdf-to-word")
async def convert_pdf_to_word(file: UploadFile = File(...), password: str = Form(None)):
    """Convert a PDF to a .docx file (supports password protection)."""
    filename = file.filename or "document.pdf"
    if not filename.lower().endswith(".pdf"):
        raise HTTPException(400, "Please upload a PDF file")
    data = await file.read()
    if not data:
        raise HTTPException(400, "Uploaded file is empty")

    import tempfile
    import fitz

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
        tmp_pdf.write(data)
        tmp_pdf_path = tmp_pdf.name

    # Check for encryption / password
    doc = fitz.open(tmp_pdf_path)
    try:
        if doc.is_encrypted:
            if doc.authenticate(password or "") == 0:
                if not password:
                    raise HTTPException(401, "password_required")
                else:
                    raise HTTPException(401, "invalid_password")
    finally:
        doc.close()

    tmp_docx_path = tmp_pdf_path.replace(".pdf", ".docx")
    try:
        from pdf2docx import Converter
    except ImportError:
        raise HTTPException(500, "pdf2docx library not installed")

    try:
        cv = Converter(tmp_pdf_path, password=password)
        cv.convert(tmp_docx_path)
        cv.close()
    except Exception as exc:
        raise HTTPException(500, f"PDF to Word conversion failed: {exc}")

    try:
        with open(tmp_docx_path, "rb") as f:
            docx_bytes = f.read()
    except Exception as exc:
        raise HTTPException(500, f"Failed to read converted file: {exc}")
    finally:
        import os as _os
        for p in [tmp_pdf_path, tmp_docx_path]:
            try:
                _os.unlink(p)
            except Exception:
                pass

    stem = re.sub(r"\.pdf$", "", filename, flags=re.I)
    return Response(
        docx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{stem}.docx"'}
    )


@app.post("/api/files/{file_id}/pdf-to-word")
def convert_stored_pdf_to_word(file_id: str, password: str = None):
    """Convert an already-uploaded PDF to a .docx file (supports password protection)."""
    meta = _meta_or_404(file_id)
    pdf_path = storage.path_for(file_id)
    docx_path = pdf_path.with_suffix(".docx")

    import fitz
    doc = fitz.open(str(pdf_path))
    try:
        if doc.is_encrypted:
            if doc.authenticate(password or "") == 0:
                if not password:
                    raise HTTPException(401, "password_required")
                else:
                    raise HTTPException(401, "invalid_password")
    finally:
        doc.close()

    try:
        from pdf2docx import Converter
    except ImportError:
        raise HTTPException(500, "pdf2docx library not installed")

    try:
        cv = Converter(str(pdf_path), password=password)
        cv.convert(str(docx_path))
        cv.close()
    except Exception as exc:
        raise HTTPException(500, f"PDF to Word conversion failed: {exc}")

    return {"ok": True}


@app.get("/api/files/{file_id}/download-word")
def download_stored_word(file_id: str):
    """Download the converted Word file for a stored PDF."""
    meta = _meta_or_404(file_id)
    docx_path = storage.path_for(file_id).with_suffix(".docx")
    if not docx_path.exists():
        raise HTTPException(404, "Word file not found. Please convert first.")

    stem = re.sub(r"\.pdf$", "", meta.get("name", "document"), flags=re.I)
    return Response(
        docx_path.read_bytes(),
        media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        headers={"Content-Disposition": f'attachment; filename="{stem}.docx"'},
    )


# ------------------------------------------------------------ PDF to PowerPoint

@app.post("/api/convert/pdf-to-ppt")
async def convert_pdf_to_ppt(file: UploadFile = File(...), password: str = Form(None)):
    """Convert a PDF to a PowerPoint .pptx file (supports password protection)."""
    filename = file.filename or "document.pdf"
    if not filename.lower().endswith(".pdf"):
        raise HTTPException(400, "Please upload a PDF file")
    data = await file.read()
    if not data:
        raise HTTPException(400, "Uploaded file is empty")

    import tempfile
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    import io

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp_pdf:
        tmp_pdf.write(data)
        tmp_pdf_path = tmp_pdf.name

    # Check for encryption / password
    doc = fitz.open(tmp_pdf_path)
    try:
        if doc.is_encrypted:
            if doc.authenticate(password or "") == 0:
                if not password:
                    raise HTTPException(401, "password_required")
                else:
                    raise HTTPException(401, "invalid_password")
        
        # Get page sizes & count
        page_count = doc.page_count
        first_page = doc[0]
        page_w = first_page.rect.width
        page_h = first_page.rect.height
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(500, f"Failed to parse PDF: {exc}")
    finally:
        doc.close()

    try:
        # Render PDF pages to high quality images
        entries = pdf_ops.export_images(Path(tmp_pdf_path), list(range(page_count)), fmt="png", dpi=150, password=password)
        if not entries:
            raise HTTPException(400, "Failed to render PDF pages")

        prs = Presentation()
        prs.slide_width = Inches(page_w / 72.0)
        prs.slide_height = Inches(page_h / 72.0)
        blank_slide_layout = prs.slide_layouts[6]

        for name, img_data in entries:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(io.BytesIO(img_data), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

        ppt_buf = io.BytesIO()
        prs.save(ppt_buf)
        pptx_bytes = ppt_buf.getvalue()
    except Exception as exc:
        raise HTTPException(500, f"PDF to PowerPoint conversion failed: {exc}")
    finally:
        import os as _os
        try:
            _os.unlink(tmp_pdf_path)
        except Exception:
            pass

    stem = re.sub(r"\.pdf$", "", filename, flags=re.I)
    return Response(
        pptx_bytes,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{stem}.pptx"'}
    )


@app.post("/api/files/{file_id}/pdf-to-ppt")
def convert_stored_pdf_to_ppt(file_id: str, password: str = None):
    """Convert an already-uploaded PDF to a .pptx file (supports password protection)."""
    meta = _meta_or_404(file_id)
    pdf_path = storage.path_for(file_id)
    pptx_path = pdf_path.with_suffix(".pptx")

    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    import io

    doc = fitz.open(str(pdf_path))
    try:
        if doc.is_encrypted:
            if doc.authenticate(password or "") == 0:
                if not password:
                    raise HTTPException(401, "password_required")
                else:
                    raise HTTPException(401, "invalid_password")
        
        # Get page sizes & count
        page_count = doc.page_count
        first_page = doc[0]
        page_w = first_page.rect.width
        page_h = first_page.rect.height
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(500, f"Failed to parse PDF: {exc}")
    finally:
        doc.close()

    try:
        # Render PDF pages to high quality images
        entries = pdf_ops.export_images(pdf_path, list(range(page_count)), fmt="png", dpi=150, password=password)
        if not entries:
            raise HTTPException(400, "Failed to render PDF pages")

        prs = Presentation()
        prs.slide_width = Inches(page_w / 72.0)
        prs.slide_height = Inches(page_h / 72.0)
        blank_slide_layout = prs.slide_layouts[6]

        for name, img_data in entries:
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(io.BytesIO(img_data), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

        prs.save(str(pptx_path))
    except Exception as exc:
        raise HTTPException(500, f"PDF to PowerPoint conversion failed: {exc}")

    return {"ok": True}


@app.get("/api/files/{file_id}/download-ppt")
def download_stored_ppt(file_id: str):
    """Download the converted PowerPoint file for a stored PDF."""
    meta = _meta_or_404(file_id)
    pptx_path = storage.path_for(file_id).with_suffix(".pptx")
    if not pptx_path.exists():
        raise HTTPException(404, "PowerPoint file not found. Please convert first.")

    stem = re.sub(r"\.pdf$", "", meta.get("name", "document"), flags=re.I)
    return Response(
        pptx_path.read_bytes(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{stem}.pptx"'},
    )


# ------------------------------------------------------------ Word to PDF

def convert_docx_to_pdf_native(data: bytes, filename: str) -> bytes:
    """Use MS Word COM interface (Windows only) to convert DOCX -> PDF."""
    import tempfile, os as _os
    import pythoncom
    import win32com.client as win32

    with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as f:
        f.write(data)
        docx_path = f.name

    pdf_path = docx_path.replace(".docx", ".pdf")
    word = None
    pythoncom.CoInitialize()
    try:
        word = win32.Dispatch("Word.Application")
        word.Visible = False
        doc = word.Documents.Open(docx_path)
        doc.SaveAs(pdf_path, FileFormat=17)  # 17 = wdFormatPDF
        doc.Close(False)
        with open(pdf_path, "rb") as f:
            return f.read()
    finally:
        for p in [docx_path, pdf_path]:
            try:
                _os.unlink(p)
            except Exception:
                pass
        try:
            if word:
                word.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()


def _convert_word_to_pdf_reportlab(data: bytes) -> bytes:
    """Fallback converter using python-docx + reportlab."""
    from docx import Document as DocxDocument
    from docx.shared import RGBColor
    import io
    try:
        docx_doc = DocxDocument(io.BytesIO(data))
    except Exception as exc:
        raise HTTPException(400, f"Could not open Word file: {exc}")

    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.pagesizes import A4
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY

    ALIGN_MAP = {"LEFT": TA_LEFT, "CENTER": TA_CENTER, "RIGHT": TA_RIGHT, "JUSTIFY": TA_JUSTIFY}

    buf = io.BytesIO()
    rldoc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=54, rightMargin=54, topMargin=54, bottomMargin=54
    )
    base_styles = getSampleStyleSheet()

    def _make_style(name, parent_name="Normal", **kw):
        base = base_styles.get(parent_name, base_styles["Normal"])
        return ParagraphStyle(name=name, parent=base, **kw)

    heading_sizes = {1: 20, 2: 16, 3: 14, 4: 12, 5: 11, 6: 10}
    story = []

    def _run_markup(run) -> str:
        txt = run.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        if not txt:
            return ""
        try:
            rgb = run.font.color.rgb if run.font.color and run.font.color.type else None
            col_hex = f"#{rgb}" if rgb else None
        except Exception:
            col_hex = None
        if col_hex:
            txt = f'<font color="{col_hex}">{txt}</font>'
        if run.bold:
            txt = f"<b>{txt}</b>"
        if run.italic:
            txt = f"<i>{txt}</i>"
        if run.underline:
            txt = f"<u>{txt}</u>"
        size = run.font.size
        if size:
            pt_size = max(7, min(72, int(size / 12700)))
            txt = f'<font size="{pt_size}">{txt}</font>'
        return txt

    for blk in docx_doc.element.body:
        tag = blk.tag.split("}")[-1] if "}" in blk.tag else blk.tag

        if tag == "p":
            from docx.text.paragraph import Paragraph as DocxParagraph
            para = DocxParagraph(blk, docx_doc)
            style_name = para.style.name if para.style else ""
            # Robustly extract alignment name from python-docx enum
            align_str = "LEFT"
            if para.alignment is not None:
                if hasattr(para.alignment, "name"):
                    align_str = para.alignment.name
                else:
                    _val_map = {0: "LEFT", 1: "CENTER", 2: "RIGHT", 3: "JUSTIFY"}
                    align_str = _val_map.get(int(para.alignment), "LEFT")
            rl_align = ALIGN_MAP.get(align_str, TA_LEFT)

            markup = "".join(_run_markup(r) for r in para.runs)
            if not markup.strip():
                story.append(Spacer(1, 6))
                continue

            if style_name.startswith("Heading"):
                try:
                    level = int(style_name.split()[-1])
                except ValueError:
                    level = 1
                font_size = heading_sizes.get(level, 13)
                sty = _make_style(
                    f"H{level}_{id(para)}", parent_name="Normal",
                    fontSize=font_size, leading=font_size * 1.3,
                    spaceBefore=8, spaceAfter=4,
                    fontName="Helvetica-Bold", alignment=rl_align
                )
            else:
                sty = _make_style(
                    f"Body_{id(para)}", parent_name="Normal",
                    fontSize=11, leading=15, spaceAfter=3,
                    alignment=rl_align
                )
            story.append(Paragraph(markup, sty))

        elif tag == "tbl":
            from docx.table import Table as DocxTable
            tbl = DocxTable(blk, docx_doc)
            table_data = []
            for row in tbl.rows:
                row_cells = []
                for cell in row.cells:
                    cell_text = " ".join(p.text for p in cell.paragraphs).strip()
                    cell_text = cell_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                    row_cells.append(Paragraph(cell_text, base_styles["Normal"]))
                table_data.append(row_cells)
            if table_data:
                col_count = max(len(r) for r in table_data)
                avail_w = A4[0] - 108
                col_w = avail_w / col_count if col_count else avail_w
                rl_table = Table(table_data, colWidths=[col_w] * col_count, repeatRows=1)
                rl_table.setStyle(TableStyle([
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f0f2f8")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                ]))
                story.append(rl_table)
                story.append(Spacer(1, 6))

    if not story:
        raise HTTPException(422, "The Word document appears to be empty or contains unsupported content only")

    try:
        rldoc.build(story)
    except Exception as exc:
        raise HTTPException(500, f"PDF rendering failed: {exc}")

    pdf_bytes = buf.getvalue()
    if not pdf_bytes:
        raise HTTPException(500, "PDF rendering produced an empty file")
    return pdf_bytes


@app.post("/api/convert/word-to-pdf")
async def convert_word_to_pdf(file: UploadFile = File(...)):
    """
    Accept a .docx/.doc file upload and convert it to PDF.
    Uses MS Word COM interface via docx2pdf if available on Windows,
    otherwise falls back to python-docx + reportlab.
    Returns a new file entry (stored PDF) so the caller can download it.
    """
    filename = file.filename or "document.docx"
    if not re.search(r"\.(docx?|DOC)$", filename):
        raise HTTPException(400, "Please upload a Word file (.doc or .docx)")

    data = await file.read()
    if not data:
        raise HTTPException(400, "Uploaded file is empty")

    pdf_bytes = None

    # Try using native MS Word conversion first if on Windows
    if sys.platform == "win32":
        try:
            from starlette.concurrency import run_in_threadpool
            pdf_bytes = await run_in_threadpool(convert_docx_to_pdf_native, data, filename)
            print("Native Word conversion succeeded.")
        except Exception as exc:
            print(f"Native Word conversion failed (falling back): {exc}")

    # Fallback path if native conversion wasn't run or failed
    if not pdf_bytes:
        pdf_bytes = _convert_word_to_pdf_reportlab(data)

    # Store as a new file entry
    base_name = re.sub(r"\.(docx?|DOC)$", "", filename)
    safe_pdf_name = _safe_name(f"{base_name}.pdf")
    try:
        page_count = pdf_ops.page_count(pdf_bytes)
    except Exception:
        page_count = 1
    meta = storage.create(safe_pdf_name, pdf_bytes, page_count)
    return {
        "ok": True,
        "file_id": meta["id"],
        "name": meta["name"],
        "pages": meta["pages"],
        "size": meta["size"],
    }


if __name__ == '__main__':
    import uvicorn
    port = int(sys.argv[2]) if len(sys.argv) > 2 and sys.argv[1] == '--port' else 8000
    uvicorn.run(app, host="127.0.0.1", port=port, log_level="info")
